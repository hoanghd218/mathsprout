#!/usr/bin/env python3
"""Generate STEMFEST 2026 MathSprout presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BRAIN_DIR = "/Users/tonyhoang/.gemini/antigravity/brain/57256b18-a749-4d7c-91ae-2fa2119e146e"
OUTPUT = os.path.join(SCRIPT_DIR, "MathSprout-STEMFEST-2026.pptx")

# Screenshots
IMG_DASHBOARD = os.path.join(BRAIN_DIR, "mathsprout_dashboard_1778407508910.png")
IMG_PRACTICE = os.path.join(BRAIN_DIR, "mathsprout_practice_1778407539227.png")
IMG_CHAT = os.path.join(BRAIN_DIR, "mathsprout_ai_chat_1778407575356.png")
IMG_SKILL = os.path.join(BRAIN_DIR, "mathsprout_skill_map_1778407613413.png")
IMG_AI_DEMO1 = os.path.join(BRAIN_DIR, "ai_demo_question1_1778408293960.png")
IMG_AI_DEMO2 = os.path.join(BRAIN_DIR, "ai_demo_question2_1778408505867.png")

# Colors
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
GREEN_MID = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xFF, 0x6F, 0x00)
BLUE = RGBColor(0x15, 0x65, 0xC0)
YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0


def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, w, h, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_text(slide, left, top, w, h, items, size=18, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_bg(slide, RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xC8, 0xE6, 0xC9))

add_text(slide, Inches(0), Inches(1.2), W, Inches(1.2),
         "🌱 MathSprout", size=60, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(2.5), W, Inches(0.8),
         "Ứng dụng học Toán có trợ lý AI cho học sinh tiểu học", size=28, color=GREEN_MID, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.5), W, Inches(0.6),
         "\"10 phút mỗi ngày — Tiến bộ lớn cả năm!\"", size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.4), Inches(2.3), Pt(3))
div.fill.solid()
div.fill.fore_color.rgb = GREEN_LIGHT
div.line.fill.background()

add_text(slide, Inches(0), Inches(4.8), W, Inches(0.5),
         "Trần Nhật Minh — Lớp 2A5, Vinschool Ocean Park 1", size=22, color=BLACK, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.5), W, Inches(0.5),
         "STEMFEST 2026", size=20, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(6.0), W, Inches(0.5),
         "Chủ đề: \"Our World, Our Responsibility\"", size=18, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, RGBColor(0xFC, 0xE4, 0xEC))

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "1️⃣  Vấn đề thực tế", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

# Left column
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0xFF, 0xEB, 0xEE))
add_text(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.5),
         "😔 Khi em học Toán...", size=24, color=BLACK, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(2.4), Inches(5.2), Inches(4),
    [
        "📚  Có nhiều bài em không hiểu được luôn",
        "🏠  Về nhà có rất nhiều bài tập khó",
        "👨‍👩‍👧  Bố mẹ rất bận đi làm, không kèm được",
        "🧑‍🏫  Cách giải của người lớn ≠ cách học sinh lớp 2",
        "👫  Nhiều bạn trong lớp cũng thấy Toán khó nhất",
    ], size=20, color=BLACK)

# Right column - quote
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), YELLOW_BG)
add_text(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
         "💬 Khảo sát bạn cùng lớp:", size=24, color=ORANGE, bold=True)
add_bullet_text(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(3.5),
    [
        "\"Khi không hiểu bài, em không biết hỏi ai\"",
        "\"Bố mẹ giải nhanh quá, em càng không hiểu\"",
        "\"Toán là môn khó nhất!\"",
        "\"Gõ bài Toán bằng bàn phím rất khó\"",
    ], size=20, color=GRAY)

add_text(slide, Inches(7.1), Inches(5.3), Inches(5.2), Inches(1),
         "→ Toán là nền tảng cho mọi môn\n→ Cần \"người bạn\" luôn sẵn sàng giúp đỡ!", size=18, color=GREEN_DARK, bold=True)


# ============================================================
# SLIDE 3: Idea / Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, GREEN_BG)

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "2️⃣  Giải pháp — MathSprout 🌱", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(1.2), W, Inches(0.6),
         "\"Tại sao mình không có một bạn trợ lý AI luôn sẵn sàng giúp các bạn nhỏ học Toán?\"",
         size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# 3 key points as cards
cards = [
    ("🤖", "Không bao giờ bận", "Bạn Sprout luôn sẵn sàng\n24/7, bất kỳ lúc nào"),
    ("💚", "Luôn kiên nhẫn", "Giải thích lại bao nhiêu\nlần cũng được"),
    ("🧒", "Nói theo cách của em", "Dùng ngôn ngữ lớp 1-3\nkhông phải người lớn"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(1.0 + i * 4.0)
    add_shape(slide, x, Inches(2.3), Inches(3.5), Inches(2.5), WHITE)
    add_text(slide, x, Inches(2.5), Inches(3.5), Inches(0.6), icon, size=44, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.2), Inches(3.5), Inches(0.5), title, size=22, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.8), Inches(3.5), Inches(1), desc, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Sprout meaning
add_shape(slide, Inches(3.0), Inches(5.2), Inches(7.3), Inches(1.5), RGBColor(0xF1, 0xF8, 0xE9))
add_text(slide, Inches(3.3), Inches(5.4), Inches(6.7), Inches(1.2),
         "🌱 Sprout = Mầm Cây\nBạn ấy giúp kiến thức của em lớn lên mỗi ngày!",
         size=22, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Prototype - Dashboard + Practice (with screenshots)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, GREEN_BG)

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "3️⃣  Sản phẩm — 3 phần chính", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

# Feature 1 - Practice
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
         "1️⃣  11 dạng bài tập Toán", size=24, color=BLUE, bold=True)
add_text(slide, Inches(0.6), Inches(2.1), Inches(2.8), Inches(2),
         "• Phép cộng, trừ, nhân, chia\n• Đo lường, hình học, phân số\n• Lớp 1 đến lớp 5\n• Song ngữ Anh-Việt 🇬🇧🇻🇳\n• Nút đọc to 🔊",
         size=16, color=BLACK)
if os.path.exists(IMG_PRACTICE):
    slide.shapes.add_picture(IMG_PRACTICE, Inches(3.4), Inches(2.0), Inches(3.0))

# Feature 2 - AI Chat
add_shape(slide, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.5), RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.6), Inches(0.5),
         "2️⃣  AI Sprout Assistant", size=24, color=GREEN_DARK, bold=True)
add_text(slide, Inches(7.1), Inches(2.1), Inches(2.5), Inches(2),
         "• Chụp ảnh đề bài 📷\n• Hoặc gõ câu hỏi\n• Giải thích từng bước\n• Đọc to Anh + Việt\n• Kiên nhẫn, không vội",
         size=16, color=BLACK)
if os.path.exists(IMG_CHAT):
    slide.shapes.add_picture(IMG_CHAT, Inches(9.6), Inches(2.0), Inches(3.0))

add_text(slide, Inches(0), Inches(6.5), W, Inches(0.6),
         "3️⃣  Trợ lý gia sư thông minh: Ghi nhớ bài sai → Gợi ý luyện đúng chỗ yếu → Tiến bộ nhanh hơn!",
         size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: Dashboard + Skill Map screenshots
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, GREEN_BG)

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "3️⃣  Demo App — Dashboard & Skill Map", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

if os.path.exists(IMG_DASHBOARD):
    slide.shapes.add_picture(IMG_DASHBOARD, Inches(0.5), Inches(1.5), Inches(5.8))
    add_text(slide, Inches(0.5), Inches(6.2), Inches(5.8), Inches(0.5),
             "🏠 Dashboard — Mầm cây, Streak, XP", size=18, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

if os.path.exists(IMG_SKILL):
    slide.shapes.add_picture(IMG_SKILL, Inches(6.8), Inches(1.5), Inches(5.8))
    add_text(slide, Inches(6.8), Inches(6.2), Inches(5.8), Inches(0.5),
             "📊 Skill Map — Bản đồ kỹ năng", size=18, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: AI Demo - Q&A with Sprout
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, RGBColor(0xE8, 0xF5, 0xE9))

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "3️⃣  Demo AI — Hỏi đáp với Sprout 🌱", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(1.1), W, Inches(0.5),
         "Em hỏi → Sprout giải thích từng bước bằng cả tiếng Anh và tiếng Việt", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# Left: Addition demo
add_shape(slide, Inches(0.4), Inches(1.8), Inches(6.2), Inches(5.2), RGBColor(0xE3, 0xF2, 0xFD))
add_text(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(0.5),
         "💬 \"How do I add 27 + 35?\"", size=22, color=BLUE, bold=True)
add_text(slide, Inches(0.6), Inches(2.5), Inches(5.8), Inches(0.8),
         "Sprout giải thích cộng hàng chục rồi hàng đơn vị,\ndùng cách đếm ngón tay — song ngữ Anh-Việt", size=16, color=GRAY)
if os.path.exists(IMG_AI_DEMO1):
    slide.shapes.add_picture(IMG_AI_DEMO1, Inches(0.8), Inches(3.3), Inches(5.5))

# Right: Fraction demo
add_shape(slide, Inches(6.9), Inches(1.8), Inches(6.0), Inches(5.2), RGBColor(0xFf, 0xF3, 0xE0))
add_text(slide, Inches(7.1), Inches(1.9), Inches(5.6), Inches(0.5),
         "💬 \"What's a fraction?\"", size=22, color=ORANGE, bold=True)
add_text(slide, Inches(7.1), Inches(2.5), Inches(5.6), Inches(0.8),
         "Sprout dùng ví dụ pizza 🍕 và bánh kem 🎂\nđể giải thích phân số — rất dễ hiểu!", size=16, color=GRAY)
if os.path.exists(IMG_AI_DEMO2):
    slide.shapes.add_picture(IMG_AI_DEMO2, Inches(7.3), Inches(3.3), Inches(5.3))


# ============================================================
# SLIDE 7: Test & Improve
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, RGBColor(0xE3, 0xF2, 0xFD))

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "4️⃣  Thử nghiệm & Cải tiến", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(1.2), W, Inches(0.5),
         "Thử app với chính em và các bạn trong lớp trong 1 tuần", size=22, color=GRAY, align=PP_ALIGN.CENTER)

# 3 improvement cards
improvements = [
    ("🔊", "Thêm nút đọc to", "Có bạn chưa đọc nhanh được\n→ Bấm nút nghe đọc to"),
    ("🇬🇧🇻🇳", "Chế độ song ngữ", "Các bạn trường quốc tế\ncũng dùng được"),
    ("📷", "Chụp ảnh đề bài", "Gõ bài Toán bằng bàn phím\nrất khó → chụp ảnh nhanh hơn"),
]
for i, (icon, title, desc) in enumerate(improvements):
    x = Inches(1.0 + i * 4.0)
    add_shape(slide, x, Inches(2.2), Inches(3.5), Inches(3.2), WHITE)
    # Number circle
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(2.4), Inches(0.8), Inches(0.8))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = GREEN_LIGHT
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x, Inches(3.4), Inches(3.5), Inches(0.5), f"{icon} {title}", size=22, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.0), Inches(3.5), Inches(1.2), desc, size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.8), W, Inches(0.8),
         "→ Lắng nghe feedback → Cải tiến sản phẩm → Sản phẩm tốt hơn cho các bạn!",
         size=20, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Impact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, WHITE, RGBColor(0xFf, 0xF3, 0xE0))

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "5️⃣  Tác động", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

impacts = [
    ("💪", "Không còn sợ Toán", "Luôn có bạn Sprout giúp đỡ"),
    ("📈", "Học giỏi hơn", "Luyện đúng vào điểm yếu"),
    ("🏠", "Tự học khi bố mẹ bận", "Không phải đợi ai cả"),
    ("🎉", "Học Toán thật vui!", "Gamification + Rewards"),
]
for i, (icon, title, desc) in enumerate(impacts):
    x = Inches(0.5 + i * 3.2)
    add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(2.8), WHITE)
    add_text(slide, x, Inches(1.7), Inches(2.8), Inches(0.7), icon, size=48, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.5), Inches(2.8), Inches(0.5), title, size=20, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.1), Inches(2.8), Inches(0.8), desc, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Slogan
add_shape(slide, Inches(2.5), Inches(4.8), Inches(8.3), Inches(1.5), GREEN_LIGHT)
add_text(slide, Inches(2.5), Inches(5.0), Inches(8.3), Inches(1.2),
         "🌟 \"10 phút mỗi ngày — Tiến bộ lớn cả năm!\" 🌟",
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8: Student Voice
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, GREEN_BG, RGBColor(0xC8, 0xE6, 0xC9))

add_text(slide, Inches(0), Inches(0.3), W, Inches(0.8),
         "6️⃣  Ý kiến học sinh", size=40, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

# Quote box
add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), WHITE)
add_text(slide, Inches(2.0), Inches(1.8), Inches(9.3), Inches(1),
         "\"Em rất tự hào về dự án này, vì đây là sản phẩm em tự nghĩ ra\ntừ chính khó khăn của em khi học Toán.\"",
         size=24, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.0), Inches(3.0), Inches(9.3), Inches(1),
         "\"Em mong rằng MathSprout sẽ giúp được thật nhiều bạn nhỏ\nở Việt Nam yêu thích môn Toán hơn.\"",
         size=22, color=BLACK, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.0), Inches(4.2), Inches(9.3), Inches(1),
         "\"Vì khi mình hiểu Toán rồi, mình sẽ thấy Toán...\nthật ra rất là vui!\" 😊",
         size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# What I did
add_text(slide, Inches(1.0), Inches(5.5), Inches(5.5), Inches(0.4),
         "✏️ Em đã tự làm:", size=18, color=GREEN_DARK, bold=True)
add_bullet_text(slide, Inches(1.0), Inches(5.9), Inches(5.5), Inches(1.5),
    ["Xác định vấn đề từ trải nghiệm thật", "Nghĩ ra ý tưởng + 4 chức năng chính",
     "Viết code, test app, quay video"], size=14, color=BLACK)

add_text(slide, Inches(7.0), Inches(5.5), Inches(5.5), Inches(0.4),
         "🤖 AI hỗ trợ em (theo quy định):", size=18, color=BLUE, bold=True)
add_bullet_text(slide, Inches(7.0), Inches(5.9), Inches(5.5), Inches(1.5),
    ["Giải thích khái niệm khó", "Hỗ trợ debug khi code lỗi",
     "AI là công cụ, không thay em tư duy"], size=14, color=BLACK)


# ============================================================
# SLIDE 9: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xA5, 0xD6, 0xA7))

add_text(slide, Inches(0), Inches(1.0), W, Inches(1),
         "🌱", size=72, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(2.2), W, Inches(1),
         "Cảm ơn ban giám khảo\nvà các thầy cô đã lắng nghe ạ!",
         size=36, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.8), W, Inches(0.6),
         "Trần Nhật Minh — Lớp 2A5 — Vinschool Ocean Park 1",
         size=24, color=BLACK, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(4.5), W, Inches(0.6),
         "STEMFEST 2026 🌱", size=22, color=GREEN_MID, bold=True, align=PP_ALIGN.CENTER)

# Quote
add_shape(slide, Inches(3.0), Inches(5.3), Inches(7.3), Inches(1.3), WHITE)
add_text(slide, Inches(3.3), Inches(5.5), Inches(6.7), Inches(1),
         "\"A small sprout today, a big tree tomorrow.\"\nMột mầm cây nhỏ hôm nay — Một cây lớn ngày mai.",
         size=20, color=GREEN_DARK, align=PP_ALIGN.CENTER)

# Save
prs.save(OUTPUT)
print(f"✅ Presentation saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
